import os
import re
from pptx import Presentation
from docx import Document
from docx.shared import Inches

def clean_text(text):
    # Remove control characters and NULL bytes
    return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', text)

def ppt_to_word(pptx_path, word_path, image_output_folder="assets/extracted_images"):
    prs = Presentation(pptx_path)
    doc = Document()

    os.makedirs(image_output_folder, exist_ok=True)

    for i, slide in enumerate(prs.slides):
        doc.add_heading(f"Slide {i + 1}", level=1)

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                text = clean_text(text)
                if text:
                    doc.add_paragraph(text)

            elif shape.shape_type == 13:  # picture
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                image_filename = f"slide_{i+1}_image.{ext}"
                image_path = os.path.join(image_output_folder, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                doc.add_picture(image_path, width=Inches(5))

        doc.add_page_break()

    doc.save(word_path)
    return word_path
