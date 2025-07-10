
import os
import io
import re
import tempfile
import subprocess
import streamlit as st
from PIL import Image
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_path
import yt_dlp
import whisper

st.set_page_config(page_title="Python Tools WebUI", layout="wide")
st.title("🧰 Python Tools Web UI")

tool = st.sidebar.selectbox("Select a tool", [
    "🧠 Whisper Transcription + Translation",
    "🖼️ PDF to JPEG (Google Slides)",
    "🔄 Flip JPEG Images Horizontally",
    "🧾 Combine Flipped Images into PDF",
    "📤 Export PPTX Text & Images to Word",
    "📽️ YouTube Downloader",
    "📥 Instagram Downloader (Public Only)",
    "🖼️ Image Format Converter"
])

# === 1. Whisper Transcription & Translation ===
if tool == "🧠 Whisper Transcription + Translation":
    st.subheader("🎧 Transcribe or Translate Audio/Video")
    audio_file = st.file_uploader("Upload MP3, WAV, or MP4", type=["mp3", "wav", "mp4"])
    task = st.selectbox("Choose task", ["transcribe", "translate"])
    src_lang = st.selectbox("Source language (or 'auto')", ["auto", "en", "hi", "mr", "ja"])
    out_format = st.radio("Download format", ["Text (.txt)", "Subtitles (.srt)"])

    if st.button("Run Whisper") and audio_file:
        with st.spinner("Processing..."):
            with tempfile.NamedTemporaryFile(delete=False, suffix=audio_file.name) as tmp:
                tmp.write(audio_file.read())
                model = whisper.load_model("base")
                result = model.transcribe(tmp.name, task=task, language=None if src_lang == "auto" else src_lang)

            if out_format == "Text (.txt)":
                text_output = result["text"]
                st.text_area("Transcript", text_output, height=300)
                st.download_button("Download TXT", text_output, file_name="transcript.txt")
            else:
                segments = result["segments"]
                srt = ""
                for i, seg in enumerate(segments):
                    srt += f"{i+1}\n"
                    start = seg["start"]
                    end = seg["end"]
                    srt += f"{int(start//60):02}:{int(start%60):02}:{int((start*1000)%1000):03} --> {int(end//60):02}:{int(end%60):02}:{int((end*1000)%1000):03}\n"
                    srt += seg["text"].strip() + "\n\n"
                st.download_button("Download SRT", srt, file_name="transcript.srt")

# === 2. PDF to JPEG ===
elif tool == "🖼️ PDF to JPEG (Google Slides)":
    st.subheader("📄 Convert PDF to JPEG")
    pdf_file = st.file_uploader("Upload exported PDF (from Google Slides)", type=["pdf"])
    if st.button("Convert to JPEG") and pdf_file:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(pdf_file.read())
            images = convert_from_path(tmp.name)
            for i, img in enumerate(images):
                buf = io.BytesIO()
                img.save(buf, format="JPEG")
                st.image(img, caption=f"Slide {i+1}")
                st.download_button(f"Download Slide {i+1}", buf.getvalue(), file_name=f"slide_{i+1}.jpg")

# === 3. Flip Images ===
elif tool == "🔄 Flip JPEG Images Horizontally":
    st.subheader("↔️ Flip JPG Images Horizontally")
    files = st.file_uploader("Upload JPG/JPEG files", accept_multiple_files=True, type=["jpg", "jpeg"])
    if st.button("Flip") and files:
        for img_file in files:
            img = Image.open(img_file)
            flipped = img.transpose(Image.FLIP_LEFT_RIGHT)
            buf = io.BytesIO()
            flipped.save(buf, format="JPEG")
            st.image(flipped, caption="Flipped Image")
            st.download_button("Download", buf.getvalue(), file_name="flipped_" + img_file.name)

# === 4. Combine Flipped Images into PDF ===
elif tool == "🧾 Combine Flipped Images into PDF":
    st.subheader("🧾 Combine Flipped JPGs to PDF")
    files = st.file_uploader("Upload Flipped JPGs", accept_multiple_files=True, type=["jpg", "jpeg"])

    def natural_sort_key(s): return [int(t) if t.isdigit() else t.lower() for t in re.split(r'(\d+)', s.name)]

    if st.button("Create PDF") and files:
        files = sorted(files, key=natural_sort_key)
        images = []
        for f in files:
            img = Image.open(f)
            if img.mode != "RGB":
                img = img.convert("RGB")
            images.append(img)
        pdf_buf = io.BytesIO()
        images[0].save(pdf_buf, save_all=True, append_images=images[1:], format="PDF")
        st.download_button("Download PDF", pdf_buf.getvalue(), file_name="combined.pdf")

# === 5. PPTX to Word ===
elif tool == "📤 Export PPTX Text & Images to Word":
    st.subheader("📤 Export PowerPoint content to Word")
    pptx_file = st.file_uploader("Upload PPTX", type=["pptx"])
    if st.button("Export to Word") and pptx_file:
        try:
            prs = Presentation(pptx_file)
            doc = Document()
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            doc.add_paragraph(text)
            buf = io.BytesIO()
            doc.save(buf)
            st.download_button("Download Word Document", buf.getvalue(), file_name="slides.docx")
        except Exception as e:
            st.error(f"❌ Word export failed: {e}")

# === 6. YouTube Downloader ===
elif tool == "📽️ YouTube Downloader":
    st.subheader("📽️ Download YouTube Video or Audio")
    url = st.text_input("Enter YouTube URL")
    quality = st.selectbox("Select Format", ["MP3", "WAV", "720p MP4", "1080p MP4", "4K MP4"])

    if st.button("Download") and url:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                ydl_opts = {"outtmpl": os.path.join(tmpdir, "%(title)s.%(ext)s")}
                if quality in ["MP3", "WAV"]:
                    ydl_opts.update({
                        "format": "bestaudio/best",
                        "postprocessors": [{
                            "key": "FFmpegExtractAudio",
                            "preferredcodec": quality.lower(),
                        }]
                    })
                else:
                    ydl_opts.update({
                        "format": "bestvideo[ext=mp4]+bestaudio[ext=m4a]/mp4",
                        "merge_output_format": "mp4"
                    })
                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    info = ydl.extract_info(url, download=True)
                    downloaded_file = ydl.prepare_filename(info).rsplit(".", 1)[0] + ".mp4"

                fixed_output = os.path.join(tmpdir, "fixed_output.mp4")
                subprocess.run([
                    "ffmpeg", "-y", "-i", downloaded_file,
                    "-r", "30",
                    "-profile:v", "high",
                    "-level", "4.0",
                    "-pix_fmt", "yuv420p",
                    "-c:v", "libx264", "-preset", "fast", "-crf", "23",
                    "-c:a", "aac",
                    "-movflags", "+faststart",
                    fixed_output
                ], check=True)
                with open(fixed_output, "rb") as f:
                    st.download_button("⬇️ Download Final MP4", f, file_name="fixed_video.mp4", mime="video/mp4")

        except Exception as e:
            st.error(f"❌ Failed to download: {e}")

# === 7. Instagram Downloader ===
elif tool == "📥 Instagram Downloader (Public Only)":
    st.subheader("📥 Download Instagram Reel or Post (Public)")
    insta_url = st.text_input("Enter Instagram Reel/Post URL")

    if st.button("Download") and insta_url:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                ydl_opts = {
                    "outtmpl": os.path.join(tmpdir, "%(title)s.%(ext)s")
                }
                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    info = ydl.extract_info(insta_url, download=True)
                    filename = ydl.prepare_filename(info)
                    with open(filename, "rb") as f:
                        st.download_button("Download File", f, file_name=os.path.basename(filename))
        except Exception as e:
            st.error(f"❌ Failed to download: {e}")

# === 8. Image Format Converter ===
elif tool == "🖼️ Image Format Converter":
    st.subheader("Convert images to another format")
    uploaded_images = st.file_uploader("Upload images", accept_multiple_files=True, type=["jpg", "jpeg", "png", "webp", "avif", "heic"])
    output_format = st.selectbox("Output format", ["JPG", "PNG", "WEBP", "HEIC", "AVIF"])

    format_map = {
        "JPG": "JPEG",
        "PNG": "PNG",
        "WEBP": "WEBP",
        "HEIC": "HEIC",
        "AVIF": "AVIF"
    }

    if st.button("Convert") and uploaded_images:
        for uploaded in uploaded_images:
            try:
                img = Image.open(uploaded)
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                ext = format_map[output_format]
                filename = os.path.splitext(uploaded.name)[0] + "." + output_format.lower()
                with tempfile.NamedTemporaryFile(delete=False, suffix="." + output_format.lower()) as tmp:
                    img.save(tmp.name, ext)
                    with open(tmp.name, "rb") as f:
                        st.download_button(f"Download {filename}", f, file_name=filename)
            except Exception as e:
                st.error(f"❌ Failed to convert {uploaded.name}: {e}")
